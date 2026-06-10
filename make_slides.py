from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color palette ──────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF5, 0xF7, 0xFA)   # slide background
CARD_BG    = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT     = RGBColor(0x1A, 0x73, 0xE8)   # Google-blue
ACCENT2    = RGBColor(0x0F, 0x9D, 0x58)   # green
ACCENT3    = RGBColor(0xF4, 0xB4, 0x00)   # amber
ACCENT4    = RGBColor(0xEA, 0x43, 0x35)   # red
DARK_TEXT  = RGBColor(0x1C, 0x1C, 0x1E)
MID_TEXT   = RGBColor(0x44, 0x47, 0x4F)
LIGHT_TEXT = RGBColor(0x80, 0x86, 0x96)
BORDER     = RGBColor(0xE0, 0xE4, 0xEC)
TAG_BG     = RGBColor(0xE8, 0xF0, 0xFE)   # light blue tag
TAG_FG     = RGBColor(0x18, 0x57, 0xC0)

# ── Helpers ────────────────────────────────────────────────
def blank_slide():
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT_BG
    return slide

def rect(slide, x, y, w, h, fill, radius=False):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox

def add_multiline(slide, lines, x, y, w, h, size, color=MID_TEXT, leading=1.15):
    """lines = list of (text, bold, color_override_or_None)"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, col) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Segoe UI"
        run.font.color.rgb = col if col else color
    return txBox

def top_bar(slide, label, tag_color=ACCENT):
    """Blue accent bar at top with slide label"""
    rect(slide, 0, 0, 13.33, 0.08, tag_color)
    add_text(slide, "CANSAT GROUND STATION  ·  COMMUNICATIONS SYSTEM",
             0.35, 0.12, 9, 0.35, 8, color=LIGHT_TEXT)
    add_text(slide, label, 10.5, 0.10, 2.5, 0.35, 8,
             color=LIGHT_TEXT, align=PP_ALIGN.RIGHT)

def slide_title(slide, number, title, subtitle=""):
    add_text(slide, f"{number:02d}", 0.45, 0.55, 0.8, 0.7, 36,
             bold=True, color=ACCENT)
    add_text(slide, title, 1.15, 0.52, 10, 0.65, 26,
             bold=True, color=DARK_TEXT)
    if subtitle:
        add_text(slide, subtitle, 1.15, 1.12, 10, 0.4, 12,
                 color=LIGHT_TEXT, italic=True)

def divider(slide, y):
    rect(slide, 0.45, y, 12.43, 0.012, BORDER)

def card(slide, x, y, w, h, title, lines, title_color=ACCENT):
    """White card with title + bullet lines"""
    shape = rect(slide, x, y, w, h, CARD_BG)
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    add_text(slide, title, x+0.18, y+0.13, w-0.3, 0.32,
             10, bold=True, color=title_color)
    rect(slide, x, y, 0.045, h, title_color)  # left accent strip
    yy = y + 0.45
    for line in lines:
        add_text(slide, f"  {line}", x+0.18, yy, w-0.3, 0.28, 9.5, color=MID_TEXT)
        yy += 0.27
    return shape

def tag(slide, text, x, y, color=ACCENT):
    shape = rect(slide, x, y, len(text)*0.1+0.25, 0.28, TAG_BG)
    shape.line.color.rgb = color
    shape.line.width = Pt(0.5)
    add_text(slide, text, x+0.1, y+0.04, len(text)*0.1+0.1, 0.22,
             8, bold=True, color=color, align=PP_ALIGN.LEFT)

def code_box(slide, x, y, w, h, code_text):
    shape = rect(slide, x, y, w, h, RGBColor(0xF1, 0xF3, 0xF4))
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    txBox = slide.shapes.add_textbox(Inches(x+0.18), Inches(y+0.12),
                                     Inches(w-0.3), Inches(h-0.2))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in code_text.split('\n'):
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(8.5)
        run.font.name = "Cascadia Code"
        run.font.color.rgb = RGBColor(0x1E, 0x1E, 0x2E)


# ════════════════════════════════════════════════════════════
# SLIDE 1 — System Architecture
# ════════════════════════════════════════════════════════════
s1 = blank_slide()
top_bar(s1, "01 / 05")
slide_title(s1, 1, "System Architecture", "How data travels from sensor to screen")
divider(s1, 1.55)

# Flow boxes
flow = [
    (0.35, "CanSat\nESP32", "9 sensors\nLoRa TX", ACCENT),
    (3.10, "GS\nESP32",     "LoRa RX\nUSB relay", ACCENT2),
    (5.85, "server.py",     "Parse · Filter\nBroadcast", ACCENT3),
    (8.60, "Dashboard\nBrowser", "Charts · Map\nExport", ACCENT4),
]

arrow_y = 3.05
for (x, title, sub, col) in flow:
    # box
    shape = rect(s1, x, 2.05, 2.35, 1.55, col)
    shape.line.fill.background()
    add_text(s1, title, x+0.12, 2.18, 2.1, 0.6, 13, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s1, sub, x+0.1, 2.78, 2.15, 0.55, 9,
             color=RGBColor(0xDD,0xEE,0xFF), align=PP_ALIGN.CENTER)

# Arrows between boxes
for ax in [2.72, 5.47, 8.22]:
    add_text(s1, "→", ax, arrow_y-0.52, 0.4, 0.4, 18, bold=True,
             color=LIGHT_TEXT, align=PP_ALIGN.CENTER)

# Protocol labels under arrows
protos = [(2.55, "LoRa 923.75 MHz\nSF11 · 125kHz BW"),
          (5.30, "USB Serial\n115200 baud"),
          (8.05, "WebSocket\nws://localhost:8765")]
for (px, pt) in protos:
    add_text(s1, pt, px, 3.65, 2.2, 0.55, 8,
             color=LIGHT_TEXT, align=PP_ALIGN.CENTER, italic=True)

# Detail cards bottom row
card(s1, 0.35, 4.45, 3.85, 2.6, "CanSat Sensors", [
    "BME280 — Temp / Humidity / Altitude",
    "ADXL345 — Acceleration XYZ",
    "GPS (UART) — Lat, Lon, Satellites",
    "PMSA003 — PM1.0 / PM2.5 / PM10",
    "INA219 — Voltage, Current, Power",
])
card(s1, 4.40, 4.45, 4.55, 2.6, "LoRa Parameters", [
    "Frequency  :  923.75 MHz",
    "Spreading Factor  :  SF11",
    "Bandwidth  :  125 kHz",
    "TX Power  :  17 dBm",
    "Time-on-Air  :  ~2.2 s / packet",
], title_color=ACCENT2)
card(s1, 9.15, 4.45, 3.85, 2.6, "Key Files", [
    "cansat_esp32.ino  — firmware",
    "gs_esp32.ino  — GS receiver",
    "server.py  — backend (869 lines)",
    "dashboard.js  — UI (945 lines)",
    "dashboard.html  — layout",
], title_color=ACCENT3)


# ════════════════════════════════════════════════════════════
# SLIDE 2 — Data Protocol
# ════════════════════════════════════════════════════════════
s2 = blank_slide()
top_bar(s2, "02 / 05")
slide_title(s2, 2, "Data Protocol — CSV Packet", "21 fields per packet · sent every 1 second")
divider(s2, 1.55)

add_text(s2, "ทุก 1 วินาที CanSat จะส่ง 1 packet เป็น CSV text ผ่านทั้ง LoRa และ USB Serial",
         0.45, 1.60, 12, 0.4, 10.5, color=MID_TEXT)

# Packet visual — colored field groups
groups = [
    ("ID & Time",     ["team_id", "time (s)", "packet_id"], ACCENT),
    ("GPS",           ["lat", "lon", "satellites"],         ACCENT2),
    ("Environment",   ["temp °C", "humidity %", "alt_baro m"], RGBColor(0x9C,0x27,0xB0)),
    ("IMU",           ["acc_x", "acc_y", "acc_z", "heading °"], ACCENT3),
    ("Air Quality",   ["PM1.0", "PM2.5", "PM10  µg/m³"],    RGBColor(0x00,0x96,0x88)),
    ("Power",         ["voltage V", "current mA", "watt W", "battery %"], ACCENT4),
    ("Status",        ["status (bitmask)"],                  DARK_TEXT),
]

gx = 0.35
gy = 2.10
max_h = 0
for (gname, fields, gcol) in groups:
    gh = 0.38 + len(fields) * 0.32
    if max_h < gh: max_h = gh
    gw = 1.72
    shape = rect(s2, gx, gy, gw, gh, gcol)
    shape.line.fill.background()
    add_text(s2, gname, gx+0.1, gy+0.06, gw-0.15, 0.3,
             8.5, bold=True, color=WHITE)
    for i, f in enumerate(fields):
        add_text(s2, f"• {f}", gx+0.1, gy+0.36+i*0.32, gw-0.15, 0.3,
                 8, color=RGBColor(0xDD,0xEE,0xFF))
    gx += gw + 0.12

# Status bitmask explanation
rect_shape = rect(s2, 0.35, 4.10, 12.63, 1.05, CARD_BG)
rect_shape.line.color.rgb = BORDER
rect_shape.line.width = Pt(0.75)
rect(s2, 0.35, 4.10, 0.045, 1.05, ACCENT3)

add_text(s2, "Status Bitmask", 0.55, 4.17, 3, 0.3, 9.5, bold=True, color=ACCENT3)
bits = [("Bit 0 = 1", "ASCENDING",  ACCENT2),
        ("Bit 1 = 2", "APOGEE",     ACCENT),
        ("Bit 2 = 4", "DEPLOYMENT", ACCENT3),
        ("Bit 3 = 8", "DESCENDING", RGBColor(0x9C,0x27,0xB0)),
        ("Bit 4 = 16","LANDED",     ACCENT4)]
bx = 0.55
for (bval, blabel, bcol) in bits:
    add_text(s2, bval,   bx,      4.50, 1.1, 0.25, 8, color=LIGHT_TEXT)
    add_text(s2, blabel, bx,      4.72, 1.1, 0.3, 9.5, bold=True, color=bcol)
    bx += 2.45

# Why CSV
rect_s = rect(s2, 0.35, 5.30, 12.63, 0.88, TAG_BG)
rect_s.line.color.rgb = ACCENT
rect_s.line.width = Pt(0.75)
add_text(s2, "ทำไมถึงใช้ CSV?", 0.55, 5.38, 2.5, 0.3, 9.5, bold=True, color=ACCENT)
add_text(s2,
         "Lightweight ไม่มี overhead  ·  ส่งได้เร็ว  ·  Debug ง่าย (อ่านเป็น text ได้เลย)  ·  "
         "ขนาดเล็กเหมาะกับ LoRa ที่มี bandwidth จำกัด",
         0.55, 5.68, 12.3, 0.4, 9.5, color=MID_TEXT)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — Backend Processing
# ════════════════════════════════════════════════════════════
s3 = blank_slide()
top_bar(s3, "03 / 05")
slide_title(s3, 3, "Backend Processing  —  server.py",
            "Parse · Validate · Filter · Derive  (869 lines, async Python)")
divider(s3, 1.55)

# 3 pipeline steps
steps = [
    ("① Parse & Validate",
     ["อ่าน Serial port แบบ async (auto-reconnect)",
      "แยก CSV → ตรวจทุก field",
      "รองรับ packet ซ้อนกัน 2 ชุด/frame (SF11 ToA > 2s)",
      "Field เสียหาย → ทิ้ง packet ทั้งใบ"],
     ACCENT),
    ("② Kalman Filter",
     ["Smooth ค่าที่กระตุกก่อนส่งไป dashboard",
      "alt_baro  Q=0.5  R=2.0  (noise ปานกลาง)",
      "temp      Q=0.02 R=0.5  (stable มาก)",
      "acc_x/y/z Q=80   R=150  (noisy มาก)",
      "voltage   Q=0.005 R=0.05 (stable มาก)"],
     ACCENT2),
    ("③ Derive Values",
     ["Speed — Haversine distance ÷ time (m/s)",
      "Vertical accel — จาก acc_y โดยตรง",
      "Pitch / Roll / Tilt — จาก accel vector",
      "สถิติ MAX/MIN/AVG/STDDEV แยกต่อ phase"],
     ACCENT3),
]

sx = 0.35
for (title, items, col) in steps:
    sh = 4.55
    sw = 4.07
    shape = rect(s3, sx, 1.70, sw, sh, CARD_BG)
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    rect(s3, sx, 1.70, sw, 0.42, col)
    add_text(s3, title, sx+0.15, 1.75, sw-0.2, 0.35, 11,
             bold=True, color=WHITE)
    yy = 2.22
    for item in items:
        # bullet
        add_text(s3, "▸", sx+0.15, yy, 0.3, 0.3, 9, bold=True, color=col)
        add_text(s3, item, sx+0.42, yy, sw-0.55, 0.35, 9.5, color=MID_TEXT)
        yy += 0.42
    sx += sw + 0.22

# Kalman formula note
rect_s = rect(s3, 0.35, 6.38, 12.63, 0.72, TAG_BG)
rect_s.line.color.rgb = ACCENT2
rect_s.line.width = Pt(0.75)
add_text(s3, "Kalman Filter คืออะไร?", 0.55, 6.44, 3, 0.28, 9, bold=True, color=ACCENT2)
add_text(s3,
         "ค่าเฉลี่ยที่ฉลาดขึ้น — ถ่วงน้ำหนักระหว่าง \"ค่าที่ทำนายไว้\" กับ \"ค่าที่วัดจริง\"  "
         "ยิ่ง R สูง = เชื่อ sensor น้อยลง  ·  ยิ่ง Q สูง = ยอมให้ค่าเปลี่ยนเร็ว",
         0.55, 6.72, 12.3, 0.3, 9.5, color=MID_TEXT)


# ════════════════════════════════════════════════════════════
# SLIDE 4 — Flight State Machine
# ════════════════════════════════════════════════════════════
s4 = blank_slide()
top_bar(s4, "04 / 05")
slide_title(s4, 4, "Flight State Machine",
            "ดาวเทียมรู้ได้ยังไงว่าตัวเองอยู่ที่ไหนในภารกิจ")
divider(s4, 1.55)

states = [
    ("🚀", "ASCENDING",   ACCENT,                      "alt > 0m\nกำลังขึ้น"),
    ("🎯", "APOGEE",      RGBColor(0x9C,0x27,0xB0),   "peak − now ≥ 3m\n& alt > 267m"),
    ("🪂", "DEPLOYMENT",  ACCENT3,                     "Servo กางร่ม\n(PWM 45°)"),
    ("⬇", "DESCENDING",  RGBColor(0x00,0x96,0x88),   "กำลังร่วงลง"),
    ("🏁", "LANDED",      ACCENT4,                     "alt < 50m\n& speed < 0.15 m/s"),
]

bx = 0.35
for i, (icon, name, col, cond) in enumerate(states):
    # State box
    shape = rect(s4, bx, 1.80, 2.35, 1.55, col)
    shape.line.fill.background()
    add_text(s4, icon, bx+0.85, 1.85, 0.8, 0.55, 22, align=PP_ALIGN.CENTER)
    add_text(s4, name, bx+0.1, 2.42, 2.15, 0.4, 10,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, cond, bx+0.1, 2.80, 2.15, 0.45, 8,
             color=RGBColor(0xDD,0xEE,0xFF), align=PP_ALIGN.CENTER)
    # arrow
    if i < 4:
        add_text(s4, "→", bx+2.37, 2.35, 0.22, 0.5, 16,
                 bold=True, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    bx += 2.57

# Detail cards
card(s4, 0.35, 3.65, 4.0, 3.0, "Apogee Detection Logic", [
    "1. Altitude ต้องเกิน 267 m ก่อน",
    "2. (peak − current) ≥ 3 m",
    "3. Failsafe: ถ้า sensor เสีย",
    "   → force APOGEE หลัง 10 วิ",
], title_color=RGBColor(0x9C,0x27,0xB0))

card(s4, 4.55, 3.65, 4.0, 3.0, "Servo Deployment", [
    "ทำงานที่ APOGEE ทันที",
    "หมุนไป 45° (กางร่มชูชีพ)",
    "ค้างไว้ 3 วินาที ก่อนเปลี่ยน state",
    "ไม่มีไฟฟ้า = ร่มไม่กาง (safety)",
], title_color=ACCENT3)

card(s4, 8.75, 3.65, 4.23, 3.0, "Landing Detection", [
    "alt < 50 m  AND",
    "แทบไม่เคลื่อนที่ (< 0.15 m/s)",
    "ทั้ง 2 เงื่อนไขต้องจริงพร้อมกัน",
    "→ หยุดส่ง packet ประหยัดแบต",
], title_color=ACCENT4)


# ════════════════════════════════════════════════════════════
# SLIDE 5 — Dashboard & Export
# ════════════════════════════════════════════════════════════
s5 = blank_slide()
top_bar(s5, "05 / 05")
slide_title(s5, 5, "Dashboard & Data Export",
            "Real-time web UI + Excel report สร้างอัตโนมัติ")
divider(s5, 1.55)

# 3 Tab cards
tabs = [
    ("📡  Monitor Tab",
     ["12 metric cards (Alt, Speed, Temp, PM2.5 …)",
      "6 real-time charts ด้วย Chart.js",
      "Flight phase bar (ASCENT → LANDED)",
      "อัปเดตทุก packet ที่รับเข้ามา"],
     ACCENT),
    ("🗺  Map Tab",
     ["Leaflet map + OpenStreetMap",
      "Flight trail (เส้นทางการบินสด)",
      "ระยะห่างจาก Ground Station (km)",
      "ปุ่ม FOLLOW / FIT / RESET TRAIL"],
     ACCENT2),
    ("📊  Analysis Tab",
     ["ตาราง MAX / MIN / AVG / STDDEV ต่อ phase",
      "3D Rocket orientation viewer",
      "Altitude vs Speed scatter plot",
      "Replay จากไฟล์ Excel (ปรับความเร็วได้)"],
     RGBColor(0x9C,0x27,0xB0)),
]

tx = 0.35
for (title, items, col) in tabs:
    tw = 3.95
    shape = rect(s5, tx, 1.70, tw, 3.70, CARD_BG)
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    rect(s5, tx, 1.70, tw, 0.40, col)
    add_text(s5, title, tx+0.15, 1.75, tw-0.2, 0.32,
             10.5, bold=True, color=WHITE)
    yy = 2.20
    for item in items:
        add_text(s5, "▸", tx+0.15, yy, 0.3, 0.28, 9, bold=True, color=col)
        add_text(s5, item, tx+0.42, yy, tw-0.55, 0.32, 9.5, color=MID_TEXT)
        yy += 0.38
    tx += tw + 0.27

# Excel export section
rect_e = rect(s5, 0.35, 5.60, 12.63, 1.60, CARD_BG)
rect_e.line.color.rgb = ACCENT3
rect_e.line.width = Pt(1.0)
rect(s5, 0.35, 5.60, 0.055, 1.60, ACCENT3)

add_text(s5, "Excel Auto-Export  —  cansat_telemetry_YYYYMMDD_HHMMSS.xlsx",
         0.55, 5.68, 9, 0.35, 10.5, bold=True, color=ACCENT3)
add_text(s5, "กด STOP บน dashboard → ไฟล์ถูกสร้างทันที ไม่ต้อง export เอง",
         0.55, 6.02, 9, 0.3, 9, color=LIGHT_TEXT, italic=True)

sheets = [
    ("📡 Telemetry", "raw + derived 26 คอลัมน์", ACCENT),
    ("📊 Statistics", "MAX/MIN/AVG แยก phase\n+ Bar charts", ACCENT2),
    ("📈 Charts", "Time-series 6 กราฟ\n(alt, speed, temp …)", ACCENT3),
]
sx2 = 0.55
for (sname, sdesc, scol) in sheets:
    shape = rect(s5, sx2, 6.35, 3.88, 0.73, scol)
    shape.line.fill.background()
    add_text(s5, sname,  sx2+0.15, 6.38, 3.5, 0.3, 9.5, bold=True, color=WHITE)
    add_text(s5, sdesc,  sx2+0.15, 6.65, 3.5, 0.35, 8.5,
             color=RGBColor(0xDD,0xEE,0xFF))
    sx2 += 4.10


# ── Save ────────────────────────────────────────────────────
out = r"c:\cansat_gs\CanSat_Communications_System.pptx"
prs.save(out)
print(f"Saved: {out}")
